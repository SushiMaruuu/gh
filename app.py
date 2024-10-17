import streamlit as st
from pptx import Presentation

# Function to clean text (replace unsupported characters, e.g., en-dash)
def clean_text(text):
    return text.replace('\u2013', '-')  # Replace en-dash with a simple dash

# Streamlit app
st.title("PowerPoint Text Extractor")

# File uploader for PPTX files only
uploaded_file = st.file_uploader("Upload a PowerPoint (PPTX)", type=["pptx"])

if uploaded_file is not None:
    st.write(f"Filename: {uploaded_file.name}")
    
    # Load the PowerPoint file
    presentation = Presentation(uploaded_file)
    all_text = []

    # Extract text from each slide
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):  # Only extract text-containing shapes
                slide_text.append(shape.text)
        
        # Clean and join the text for each slide
        cleaned_slide_text = clean_text("\n".join(slide_text))
        all_text.append(f"**Slide {i + 1}:**\n{cleaned_slide_text}\n")

    # Display the extracted text from all slides
    st.subheader("Extracted Text from PowerPoint")
    for slide_text in all_text:
        st.write(slide_text)

    # Optionally, show a message if no text is found in the slides
    if not any(slide_text for slide_text in all_text):
        st.write("No text found in the PowerPoint.")
